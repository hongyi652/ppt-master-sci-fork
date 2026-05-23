#!/usr/bin/env python3
"""Basic direct PPTX generator for the local web GUI."""

from __future__ import annotations

import importlib.util
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


CANVAS_SIZES = {
    "ppt169": (13.333, 7.5),
    "ppt43": (10.0, 7.5),
    "xhs": (7.5, 10.0),
    "xiaohongshu": (7.5, 10.0),
    "story": (7.5, 13.333),
}
MAX_SLIDES = 14
MAX_OUTLINE_ITEMS = 8
MAX_BULLETS_PER_SLIDE = 5
MIN_AUTO_ASSET_SCORE = 2.0
MIN_CAPTION_BIND_TOKENS = 2
DIAGNOSTIC_CANDIDATE_LIMIT = 6
FONT_FAMILY = "Microsoft YaHei"
SUPPORTED_LAYOUTS = {"insight", "cards", "comparison", "timeline", "image_focus", "spotlight", "summary"}
SUPPORTED_PAGE_RHYTHMS = {"anchor", "dense", "breathing"}
ASSET_TOKEN_PATTERN = re.compile(r"[A-Za-z0-9_+\-]{3,}|[\u4e00-\u9fff]{2,}")
FIGURE_REF_PATTERN = re.compile(r"(?:\bfig(?:ure)?\.?\s*|图\s*)([A-Za-z]?\d+(?:\.\d+)?[A-Za-z]?)", re.I)
EMU_PER_PIXEL = 9525
ASSET_METADATA_TOKEN_FIELDS = (
    "purpose",
    "search_query",
    "source_title",
    "source_author",
    "provider",
    "license_name",
    "license_tier",
    "source_kind",
    "original_filename",
    "source_namespace",
    "figure_label",
    "figure_ref_key",
    "figure_caption",
    "section_heading",
    "section_heading_path",
)
ASSET_PRIORITY_METADATA_FIELDS = (
    "purpose",
    "search_query",
    "source_title",
    "figure_caption",
    "section_heading",
)


def _load_svg_media_module():
    module_path = (
        Path(__file__).resolve().parent
        / "skills"
        / "ppt-master"
        / "scripts"
        / "svg_to_pptx"
        / "pptx_media.py"
    )
    if not module_path.exists():
        return None

    spec = importlib.util.spec_from_file_location("ppt_master_simple_builder_svg_media", module_path)
    if spec is None or spec.loader is None:
        return None

    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


_SVG_MEDIA_MODULE = _load_svg_media_module()
_CONVERT_SVG_TO_PNG_CACHED = (
    getattr(_SVG_MEDIA_MODULE, "convert_svg_to_png_cached", None) if _SVG_MEDIA_MODULE else None
)


@dataclass
class Section:
    title: str
    bullets: list[str]
    image_paths: list[Path] = field(default_factory=list)
    layout_hint: str = "insight"
    page_rhythm: str = "anchor"
    lead: str = ""


def build_presentation(
    source_text: str,
    project_name: str,
    output_path: Path,
    canvas_format: str = "ppt169",
    deck_spec: dict[str, Any] | None = None,
    source_images: list[dict[str, Any]] | None = None,
) -> dict[str, object]:
    """Build a simple editable PPTX from source text."""
    normalized = _normalize_text(source_text)
    if not normalized:
        raise ValueError("Source text is empty.")

    prs = Presentation()
    width, height = CANVAS_SIZES.get(canvas_format, CANVAS_SIZES["ppt169"])
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    title = _extract_title(normalized, project_name)
    subtitle = _extract_subtitle(normalized, title)
    sections = _parse_sections(normalized, fallback_title=title)

    if deck_spec:
        planned_title = _clean_inline(str(deck_spec.get("title", "")))
        planned_subtitle = _clean_inline(str(deck_spec.get("subtitle", "")))
        planned_sections = _sections_from_deck_spec(deck_spec, source_images or [])
        if planned_title:
            title = planned_title[:80]
        if planned_subtitle:
            subtitle = planned_subtitle[:120]
        if planned_sections:
            sections = planned_sections

    sections, asset_match_diagnostics = _apply_source_assets_to_sections(sections, source_images or [])

    _add_cover_slide(prs, title, subtitle, project_name)

    if len(sections) > 1:
        _add_outline_slide(prs, title, sections)

    slide_count = len(prs.slides)
    placed_image_count = 0
    for section in sections:
        if slide_count >= MAX_SLIDES:
            break
        chunks = _chunk(section.bullets or ["待补充内容"], MAX_BULLETS_PER_SLIDE)
        for chunk_index, chunk in enumerate(chunks):
            if slide_count >= MAX_SLIDES:
                break
            slide_title = section.title if chunk_index == 0 else f"{section.title}（续）"
            image_paths = section.image_paths if chunk_index == 0 else []
            placed_image_count += _add_content_slide(
                prs,
                slide_title,
                chunk,
                project_name,
                image_paths,
                layout_hint=section.layout_hint if chunk_index == 0 else "insight",
                page_rhythm=section.page_rhythm,
                lead=section.lead if chunk_index == 0 else (chunk[0] if chunk else section.title),
            )
            slide_count += 1

    _add_closing_slide(prs, project_name)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return {
        "title": title,
        "slide_count": len(prs.slides),
        "output_path": str(output_path),
        "placed_image_count": placed_image_count,
        "asset_match_diagnostics": asset_match_diagnostics,
    }


def _normalize_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"<!--.*?-->", "", text, flags=re.S)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_title(text: str, fallback: str) -> str:
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        match = re.match(r"^#\s+(.+)$", stripped)
        if match:
            return _clean_inline(match.group(1))[:80] or fallback
        if len(stripped) <= 48:
            return _clean_inline(stripped)[:80] or fallback
        break
    return fallback


def _extract_subtitle(text: str, title: str) -> str:
    for line in text.splitlines():
        stripped = _clean_inline(line)
        if not stripped or stripped == title:
            continue
        if stripped.startswith("#"):
            continue
        return stripped[:120]
    return "本地 GUI 直出模式生成"


def _parse_sections(text: str, fallback_title: str) -> list[Section]:
    sections: list[Section] = []
    current_title = ""
    current_lines: list[str] = []
    preface: list[str] = []

    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()
        if not stripped:
            if current_lines and current_lines[-1] != "":
                current_lines.append("")
            continue

        heading = re.match(r"^(#{1,2})\s+(.+)$", stripped)
        if heading:
            if current_title or current_lines:
                sections.append(Section(current_title or "内容概览", _lines_to_bullets(current_lines)))
            current_title = _clean_inline(heading.group(2)) or "内容概览"
            current_lines = []
            continue

        if not current_title and not sections:
            preface.append(stripped)
        else:
            current_lines.append(stripped)

    if current_title or current_lines:
        sections.append(Section(current_title or "内容概览", _lines_to_bullets(current_lines)))

    cleaned_sections = [
        Section(section.title, _dedupe_items(section.bullets))
        for section in sections
        if section.bullets
    ]

    if cleaned_sections:
        return cleaned_sections

    fallback_bullets = _lines_to_bullets(preface or text.splitlines())
    if not fallback_bullets:
        fallback_bullets = [fallback_title]
    return [Section("内容概览", _dedupe_items(fallback_bullets), lead=fallback_bullets[0])]


def _sections_from_deck_spec(deck_spec: dict[str, Any], source_images: list[dict[str, Any]]) -> list[Section]:
    raw_slides = deck_spec.get("slides") or deck_spec.get("sections") or []
    sections: list[Section] = []
    image_lookup: dict[str, Path] = {}
    ordered_images: list[Path] = []
    for raw_image in source_images:
        image_id = str(raw_image.get("id", "")).strip()
        image_path = Path(str(raw_image.get("path", "")))
        if not image_id or not image_path.exists():
            continue
        image_lookup[image_id] = image_path
        ordered_images.append(image_path)

    used_images: set[Path] = set()
    for index, raw_slide in enumerate(raw_slides, start=1):
        if not isinstance(raw_slide, dict):
            continue
        title = _clean_inline(str(raw_slide.get("title") or f"内容页 {index}"))
        bullets = raw_slide.get("bullets") or raw_slide.get("points") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        cleaned_bullets = _dedupe_items(
            [
                _clean_inline(str(item))[:180]
                for item in bullets
                if _clean_inline(str(item))
            ]
        )
        image_refs = raw_slide.get("image_refs") or raw_slide.get("images") or []
        if isinstance(image_refs, str):
            image_refs = [image_refs]
        image_paths: list[Path] = []
        for raw_ref in image_refs:
            image_path = image_lookup.get(str(raw_ref).strip())
            if image_path is None or image_path in used_images:
                continue
            image_paths.append(image_path)
            used_images.add(image_path)
            if len(image_paths) >= 2:
                break
        layout_hint = _normalize_layout_hint(str(raw_slide.get("layout_hint") or "insight"))
        page_rhythm = _normalize_page_rhythm(str(raw_slide.get("page_rhythm") or "anchor"))
        lead = _clean_inline(str(raw_slide.get("lead") or ""))[:120]
        if title and cleaned_bullets:
            sections.append(
                Section(
                    title[:60],
                    cleaned_bullets,
                    image_paths=image_paths,
                    layout_hint=layout_hint,
                    page_rhythm=page_rhythm,
                    lead=lead or cleaned_bullets[0],
                )
            )

    return sections


def _text_tokens(text: str) -> set[str]:
    tokens: set[str] = set()
    for raw_token in ASSET_TOKEN_PATTERN.findall(text or ""):
        token = raw_token.strip().lower()
        if len(token) < 2 or token.isdigit():
            continue
        tokens.add(token)
    return tokens


def _renderable_source_assets(source_images: list[dict[str, Any]]) -> list[dict[str, Any]]:
    assets: list[dict[str, Any]] = []
    for raw_asset in source_images:
        raw_path = str(raw_asset.get("path", "") or "").strip()
        if not raw_path:
            continue
        asset_path = Path(raw_path)
        if not asset_path.exists():
            continue
        asset_copy = dict(raw_asset)
        asset_copy["_path"] = asset_path
        asset_copy["_source_asset"] = raw_asset
        assets.append(asset_copy)
    return assets


def _set_asset_planning_flags(
    raw_asset: dict[str, Any],
    *,
    candidate_material: bool | None = None,
    selected_for_deck: bool | None = None,
) -> None:
    targets = [raw_asset]
    source_asset = raw_asset.get("_source_asset")
    if isinstance(source_asset, dict) and source_asset is not raw_asset:
        targets.append(source_asset)

    for target in targets:
        if candidate_material is not None:
            target["candidate_material"] = candidate_material
        if selected_for_deck is not None:
            target["selected_for_deck"] = selected_for_deck


def _asset_tokens(raw_asset: dict[str, Any]) -> set[str]:
    tokens: set[str] = set()
    for raw_tag in raw_asset.get("tags") or []:
        tokens.update(_text_tokens(str(raw_tag)))
    for key in ("alt", "context", "reference", "source_file", "latex", *ASSET_METADATA_TOKEN_FIELDS):
        tokens.update(_text_tokens(str(raw_asset.get(key) or "")))
    path_value = raw_asset.get("_path")
    if isinstance(path_value, Path):
        tokens.update(_text_tokens(path_value.stem))
    return tokens


def _priority_metadata_tokens(raw_asset: dict[str, Any]) -> set[str]:
    tokens: set[str] = set()
    for key in ASSET_PRIORITY_METADATA_FIELDS:
        tokens.update(_text_tokens(str(raw_asset.get(key) or "")))
    return tokens


def _asset_usage_count(raw_asset: dict[str, Any]) -> int:
    raw_value = raw_asset.get("usage_count")
    try:
        return max(int(raw_value), 0)
    except (TypeError, ValueError):
        return 0


def _normalize_figure_ref(text: str) -> str:
    cleaned = str(text or "").strip()
    if not cleaned:
        return ""
    cleaned = re.sub(r"(?i)\bfigure\b|\bfig\.?\b|图", "", cleaned)
    cleaned = re.sub(r"[^A-Za-z0-9.]+", "", cleaned).lower()
    if not cleaned:
        return ""
    return f"fig{cleaned}"


def _extract_figure_refs(text: str) -> set[str]:
    refs: set[str] = set()
    for match in FIGURE_REF_PATTERN.finditer(text or ""):
        ref_key = _normalize_figure_ref(match.group(1))
        if ref_key:
            refs.add(ref_key)
    return refs


def _asset_figure_refs(raw_asset: dict[str, Any]) -> set[str]:
    refs: set[str] = set()
    for key in ("figure_ref_key", "figure_label"):
        ref_key = _normalize_figure_ref(str(raw_asset.get(key) or ""))
        if ref_key:
            refs.add(ref_key)
    return refs


def _is_original_source_visual(raw_asset: dict[str, Any]) -> bool:
    source_kind = str(raw_asset.get("source_kind") or "").lower()
    if source_kind not in {"mineru", "pptx_picture", "markdown_asset"}:
        return False
    if str(raw_asset.get("figure_ref_key") or "").strip() or str(raw_asset.get("figure_caption") or "").strip():
        return True
    hint = " ".join(
        [
            str(raw_asset.get("alt") or ""),
            str(raw_asset.get("context") or ""),
            str(raw_asset.get("reference") or ""),
        ]
    ).lower()
    return bool(re.search(r"原图|原文|论文|截图|screenshot|paper|figure|fig\.?|图\s*\d", hint))


def _asset_priority(raw_asset: dict[str, Any]) -> float:
    asset_type = str(raw_asset.get("asset_type") or "image")
    if asset_type == "formula":
        return 7.0
    if asset_type in {"chart", "diagram"}:
        return 5.5
    if _is_original_source_visual(raw_asset):
        return 4.5
    return 2.0


def _section_text(section: Section) -> str:
    return " ".join([section.title, section.lead, *section.bullets])


def _section_asset_token_matches(section: Section, raw_asset: dict[str, Any]) -> tuple[set[str], set[str]]:
    section_tokens = _text_tokens(_section_text(section))
    return section_tokens & _asset_tokens(raw_asset), section_tokens & _priority_metadata_tokens(raw_asset)


def _analyze_asset_for_section(section: Section, raw_asset: dict[str, Any]) -> dict[str, Any]:
    section_text = _section_text(section)
    section_tokens = _text_tokens(section_text)
    shared_tokens, shared_priority_tokens = _section_asset_token_matches(section, raw_asset)
    shared_tokens_list = sorted(shared_tokens)
    shared_priority_tokens_list = sorted(shared_priority_tokens)
    caption_tokens = _text_tokens(str(raw_asset.get("figure_caption") or ""))
    shared_caption_tokens = sorted(section_tokens & caption_tokens)
    section_figure_refs = _extract_figure_refs(section_text)
    asset_figure_refs = _asset_figure_refs(raw_asset)
    matched_figure_refs = sorted(section_figure_refs & asset_figure_refs)
    asset_type = str(raw_asset.get("asset_type") or "image")
    source_kind = str(raw_asset.get("source_kind") or "").lower()
    usage_count = _asset_usage_count(raw_asset)
    lowered_text = section_text.lower()

    reasons: list[str] = []
    score = float(len(shared_tokens_list) * 4)
    score += float(len(shared_priority_tokens_list) * 2)

    explicit_caption_binding = _is_original_source_visual(raw_asset) and len(shared_caption_tokens) >= MIN_CAPTION_BIND_TOKENS
    explicit_binding = False
    if matched_figure_refs:
        explicit_binding = True
        score += 24.0
        reasons.append(f"explicit figure ref match: {', '.join(matched_figure_refs[:3])}")
    elif explicit_caption_binding:
        explicit_binding = True
        score += float(min(len(shared_caption_tokens), 4) * 4)
        reasons.append(f"explicit figure caption match: {', '.join(shared_caption_tokens[:4])}")
    elif shared_caption_tokens:
        score += float(min(len(shared_caption_tokens), 4) * 1.5)
        reasons.append(f"shared figure caption tokens: {', '.join(shared_caption_tokens[:4])}")

    if shared_tokens_list:
        reasons.append(f"shared asset tokens: {', '.join(shared_tokens_list[:4])}")
    if shared_priority_tokens_list:
        reasons.append(f"shared priority metadata: {', '.join(shared_priority_tokens_list[:4])}")

    if asset_type == "formula" and re.search(r"公式|方程|推导|模型|equation|formula|latex", lowered_text):
        score += 4.0
        reasons.append("formula cue matched")
    if asset_type in {"chart", "diagram"} and re.search(r"图|表|趋势|对比|流程|架构|figure|chart|graph|diagram|workflow|trend", lowered_text):
        score += 3.0
        reasons.append("chart/diagram cue matched")
    if asset_type == "image" and re.search(r"照片|图像|实验|装置|现场|photo|image|figure", lowered_text):
        score += 2.0
        reasons.append("image cue matched")
    if usage_count > 1:
        usage_bonus = min(usage_count - 1, 3) * 0.5
        score += usage_bonus
        reasons.append(f"reuse evidence bonus: {usage_bonus:.1f}")
    if source_kind in {"mineru", "pptx_picture", "markdown_asset"} and asset_type in {"chart", "diagram", "formula"}:
        score += 1.0
        reasons.append("scientific visual source bonus")
    if _is_original_source_visual(raw_asset):
        score += 1.5
        reasons.append("original source visual priority")
    if source_kind in {"mineru", "pptx_picture"} and re.search(r"原图|原文|论文|截图|screenshot|paper|figure", lowered_text):
        score += 1.0
        reasons.append("paper screenshot cue matched")

    priority = _asset_priority(raw_asset)
    score += priority * 0.5
    reasons.append(f"asset priority weight: {priority:.1f}")

    meaningful_match = bool(shared_tokens_list or shared_priority_tokens_list or matched_figure_refs or explicit_caption_binding)
    if not meaningful_match and asset_type == "formula" and re.search(r"公式|方程|推导|模型|equation|formula|latex", lowered_text):
        meaningful_match = True

    return {
        "score": score,
        "explicit_binding": explicit_binding,
        "meaningful_match": meaningful_match,
        "shared_tokens": shared_tokens_list,
        "shared_priority_tokens": shared_priority_tokens_list,
        "shared_caption_tokens": shared_caption_tokens,
        "section_figure_refs": sorted(section_figure_refs),
        "asset_figure_refs": sorted(asset_figure_refs),
        "matched_figure_refs": matched_figure_refs,
        "reasons": reasons,
    }


def _score_asset_for_section(section: Section, raw_asset: dict[str, Any]) -> float:
    return float(_analyze_asset_for_section(section, raw_asset)["score"])


def _asset_display_label(raw_asset: dict[str, Any]) -> str:
    asset_path = raw_asset.get("_path")
    fallback = asset_path.name if isinstance(asset_path, Path) else "asset"
    return _clean_inline(str(raw_asset.get("alt") or fallback))[:120] or fallback


def _serialize_asset_candidate(raw_asset: dict[str, Any], analysis: dict[str, Any], decision: str) -> dict[str, Any]:
    asset_path = raw_asset.get("_path")
    return {
        "asset_id": str(raw_asset.get("id") or ""),
        "asset_label": _asset_display_label(raw_asset),
        "asset_type": str(raw_asset.get("asset_type") or "image"),
        "source_kind": str(raw_asset.get("source_kind") or ""),
        "path": str(asset_path) if isinstance(asset_path, Path) else "",
        "score": round(float(analysis.get("score") or 0.0), 2),
        "decision": decision,
        "explicit_binding": bool(analysis.get("explicit_binding")),
        "figure_label": str(raw_asset.get("figure_label") or ""),
        "figure_caption": str(raw_asset.get("figure_caption") or "")[:180],
        "section_heading": str(raw_asset.get("section_heading") or ""),
        "matched_figure_refs": list(analysis.get("matched_figure_refs") or []),
        "shared_tokens": list(analysis.get("shared_tokens") or [])[:8],
        "shared_priority_tokens": list(analysis.get("shared_priority_tokens") or [])[:8],
        "shared_caption_tokens": list(analysis.get("shared_caption_tokens") or [])[:8],
        "reasons": list(analysis.get("reasons") or [])[:6],
    }


def _section_asset_capacity(section: Section) -> int:
    return 2 if len(section.bullets) <= 4 else 1


def _apply_source_assets_to_sections(
    sections: list[Section],
    source_images: list[dict[str, Any]],
) -> tuple[list[Section], list[dict[str, Any]]]:
    assets = _renderable_source_assets(source_images)
    diagnostics: list[dict[str, Any]] = []
    for raw_asset in assets:
        _set_asset_planning_flags(raw_asset, candidate_material=False, selected_for_deck=False)
    if not sections or not assets:
        return sections, diagnostics

    assets_by_path = {
        raw_asset.get("_path"): raw_asset
        for raw_asset in assets
        if isinstance(raw_asset.get("_path"), Path)
    }

    used_paths = {
        image_path
        for section in sections
        for image_path in section.image_paths
        if image_path.exists()
    }

    for section in sections:
        section_figure_refs = sorted(_extract_figure_refs(_section_text(section)))
        prebound_assets = [
            assets_by_path[image_path]
            for image_path in section.image_paths
            if image_path in assets_by_path
        ]
        for raw_asset in prebound_assets:
            _set_asset_planning_flags(raw_asset, candidate_material=True, selected_for_deck=True)
        capacity = max(_section_asset_capacity(section) - len(section.image_paths), 0)
        scored_assets: list[tuple[dict[str, Any], int, dict[str, Any]]] = []
        for index, raw_asset in enumerate(assets):
            asset_path = raw_asset.get("_path")
            if not isinstance(asset_path, Path) or asset_path in used_paths:
                continue
            analysis = _analyze_asset_for_section(section, raw_asset)
            if analysis["meaningful_match"] and (analysis["explicit_binding"] or analysis["score"] >= MIN_AUTO_ASSET_SCORE):
                _set_asset_planning_flags(raw_asset, candidate_material=True)
            scored_assets.append((analysis, index, raw_asset))
        scored_assets.sort(
            key=lambda item: (item[0]["explicit_binding"], item[0]["score"], -item[1]),
            reverse=True,
        )

        matched_assets: list[dict[str, Any]] = []
        selected_paths = {
            image_path
            for image_path in section.image_paths
            if isinstance(image_path, Path)
        }
        for analysis, _index, raw_asset in scored_assets:
            asset_path = raw_asset.get("_path")
            if not isinstance(asset_path, Path):
                continue
            if not analysis["explicit_binding"] and analysis["score"] < MIN_AUTO_ASSET_SCORE:
                break
            if not analysis["meaningful_match"]:
                continue
            if len(matched_assets) >= capacity:
                continue
            section.image_paths.append(asset_path)
            used_paths.add(asset_path)
            selected_paths.add(asset_path)
            matched_assets.append(raw_asset)
            _set_asset_planning_flags(raw_asset, candidate_material=True, selected_for_deck=True)

        if (matched_assets or prebound_assets) and section.layout_hint == "insight":
            if len(section.image_paths) >= 2 or any(
                str(raw_asset.get("asset_type") or "") in {"formula", "chart", "diagram"}
                for raw_asset in [*prebound_assets, *matched_assets]
            ):
                section.layout_hint = "image_focus"

        selected_assets = [
            _serialize_asset_candidate(
                raw_asset,
                _analyze_asset_for_section(section, raw_asset),
                "prebound" if raw_asset in prebound_assets else "selected",
            )
            for raw_asset in [*prebound_assets, *matched_assets]
        ]
        top_candidates: list[dict[str, Any]] = []
        for analysis, _index, raw_asset in scored_assets[:DIAGNOSTIC_CANDIDATE_LIMIT]:
            asset_path = raw_asset.get("_path")
            if not isinstance(asset_path, Path):
                continue
            if asset_path in selected_paths:
                decision = "selected"
            elif not analysis["explicit_binding"] and analysis["score"] < MIN_AUTO_ASSET_SCORE:
                decision = "below_threshold"
            elif not analysis["meaningful_match"]:
                decision = "no_meaningful_match"
            else:
                decision = "not_selected"
            top_candidates.append(_serialize_asset_candidate(raw_asset, analysis, decision))

        diagnostics.append(
            {
                "section_title": section.title,
                "lead": section.lead,
                "layout_hint": section.layout_hint,
                "capacity": _section_asset_capacity(section),
                "available_capacity": capacity,
                "section_figure_refs": section_figure_refs,
                "selected_assets": selected_assets,
                "top_candidates": top_candidates,
            }
        )

    return sections, diagnostics


def _lines_to_bullets(lines: list[str]) -> list[str]:
    bullets: list[str] = []
    paragraph_buffer: list[str] = []

    def flush_paragraph() -> None:
        if not paragraph_buffer:
            return
        paragraph = _clean_inline(" ".join(paragraph_buffer))
        paragraph_buffer.clear()
        bullets.extend(_paragraph_to_bullets(paragraph))

    for line in lines:
        stripped = line.strip()
        if not stripped:
            flush_paragraph()
            continue

        list_match = re.match(r"^\s*(?:[-*+]\s+|\d+[.)]\s+)(.+)$", stripped)
        if list_match:
            flush_paragraph()
            bullets.append(_clean_inline(list_match.group(1)))
            continue

        sub_heading = re.match(r"^#{3,6}\s+(.+)$", stripped)
        if sub_heading:
            flush_paragraph()
            bullets.append(_clean_inline(sub_heading.group(1)))
            continue

        if stripped.startswith("|") and stripped.endswith("|"):
            flush_paragraph()
            cells = [_clean_inline(cell) for cell in stripped.strip("|").split("|")]
            cells = [cell for cell in cells if cell and not set(cell) <= {"-", ":"}]
            if cells:
                bullets.append(" / ".join(cells[:4]))
            continue

        paragraph_buffer.append(stripped)

    flush_paragraph()
    return [bullet for bullet in bullets if bullet]


def _paragraph_to_bullets(paragraph: str) -> list[str]:
    if not paragraph:
        return []
    parts = [
        _clean_inline(part)
        for part in re.split(r"(?<=[。！？；.!?;])\s*", paragraph)
        if _clean_inline(part)
    ]
    if not parts:
        return []
    if len(parts) == 1 and len(parts[0]) > 90:
        parts = [
            _clean_inline(part)
            for part in re.split(r"(?<=[，、,:：])\s*", parts[0])
            if _clean_inline(part)
        ]
    return parts[:10]


def _clean_inline(text: str) -> str:
    text = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"<[^>]+>", "", text)
    text = text.replace("`", "")
    text = re.sub(r"^[>#\-\*\+\s]+", "", text)
    text = re.sub(r"[*_~]+", "", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _dedupe_items(items: list[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for item in items:
        normalized = item.lower()
        if not item or normalized in seen:
            continue
        seen.add(normalized)
        result.append(item[:180])
    return result


def _normalize_layout_hint(layout_hint: str) -> str:
    cleaned = _clean_inline(layout_hint).lower() or "insight"
    return cleaned if cleaned in SUPPORTED_LAYOUTS else "insight"


def _normalize_page_rhythm(page_rhythm: str) -> str:
    cleaned = _clean_inline(page_rhythm).lower() or "anchor"
    return cleaned if cleaned in SUPPORTED_PAGE_RHYTHMS else "anchor"


def _chunk(items: list[str], size: int) -> list[list[str]]:
    return [items[index:index + size] for index in range(0, len(items), size)] or [[]]


def _rgb(hex_value: str) -> RGBColor:
    cleaned = hex_value.lstrip("#")
    return RGBColor(int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16))


def _add_background(slide, width: int, height: int, color: str) -> None:
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    background.fill.solid()
    background.fill.fore_color.rgb = _rgb(color)
    background.line.fill.background()


def _add_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    font_size: int,
    color: str,
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_FAMILY
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _add_bullet_box(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    bullets: list[str],
    *,
    font_size: int,
    color: str,
    prefix: str = "• ",
) -> None:
    text_box = slide.shapes.add_textbox(left, top, width, height)
    frame = text_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"{prefix}{item}" if prefix else item
        paragraph.font.name = FONT_FAMILY
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = _rgb(color)
        paragraph.space_after = Pt(8)
        paragraph.line_spacing = 1.15


def _add_lead_band(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    font_size: int,
    fill_color: str = "DBEAFE",
    text_color: str = "1E3A8A",
) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    band.fill.solid()
    band.fill.fore_color.rgb = _rgb(fill_color)
    band.line.fill.background()
    _add_textbox(
        slide,
        left + int(width * 0.06),
        top + int(height * 0.18),
        int(width * 0.88),
        int(height * 0.64),
        text,
        font_size=font_size,
        color=text_color,
        bold=True,
    )


def _resolve_layout_hint(layout_hint: str, image_paths: list[Path] | None, bullets: list[str]) -> str:
    normalized = _normalize_layout_hint(layout_hint)
    if normalized == "summary":
        return "spotlight"
    if normalized == "image_focus" and not image_paths:
        return "spotlight"
    if normalized == "timeline" and len(bullets) < 3:
        return "insight"
    if normalized == "cards" and len(bullets) > 4:
        return "comparison"
    return normalized


def _rhythm_profile(page_rhythm: str, height: int) -> dict[str, float]:
    compact = height <= Inches(8)
    profiles = {
        "dense": {
            "title_size": 21 if compact else 22,
            "body_size": 16,
            "lead_size": 17,
            "gap_ratio": 0.022,
            "lead_ratio": 0.15,
        },
        "anchor": {
            "title_size": 22 if compact else 23,
            "body_size": 18,
            "lead_size": 18,
            "gap_ratio": 0.028,
            "lead_ratio": 0.17,
        },
        "breathing": {
            "title_size": 23 if compact else 24,
            "body_size": 17,
            "lead_size": 20,
            "gap_ratio": 0.034,
            "lead_ratio": 0.20,
        },
    }
    return dict(profiles[_normalize_page_rhythm(page_rhythm)])


def _add_slide_decor(slide, width: int, height: int, page_rhythm: str) -> None:
    if page_rhythm == "breathing":
        orb = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(width * 0.79), int(height * 0.12), int(width * 0.15), int(height * 0.22))
        orb.fill.solid()
        orb.fill.fore_color.rgb = _rgb("DBEAFE")
        orb.line.fill.background()
    elif page_rhythm == "dense":
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(width * 0.94), int(height * 0.12), int(width * 0.012), int(height * 0.68))
        rail.fill.solid()
        rail.fill.fore_color.rgb = _rgb("BFDBFE")
        rail.line.fill.background()


def _add_cover_slide(prs: Presentation, title: str, subtitle: str, project_name: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    width = prs.slide_width
    height = prs.slide_height
    _add_background(slide, width, height, "0F172A")

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(width * 0.035), height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb("38BDF8")
    accent.line.fill.background()

    _add_textbox(
        slide,
        int(width * 0.10),
        int(height * 0.18),
        int(width * 0.78),
        int(height * 0.26),
        title,
        font_size=28 if height < Inches(8) else 30,
        color="F8FAFC",
        bold=True,
    )
    _add_textbox(
        slide,
        int(width * 0.10),
        int(height * 0.48),
        int(width * 0.72),
        int(height * 0.14),
        subtitle,
        font_size=16,
        color="CBD5E1",
    )
    _add_textbox(
        slide,
        int(width * 0.10),
        int(height * 0.80),
        int(width * 0.70),
        int(height * 0.08),
        f"{project_name}  |  中文智能生成",
        font_size=11,
        color="94A3B8",
    )


def _add_outline_slide(prs: Presentation, title: str, sections: list[Section]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    width = prs.slide_width
    height = prs.slide_height
    _add_background(slide, width, height, "F8FAFC")

    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, int(height * 0.13))
    banner.fill.solid()
    banner.fill.fore_color.rgb = _rgb("DBEAFE")
    banner.line.fill.background()

    _add_textbox(slide, int(width * 0.07), int(height * 0.05), int(width * 0.60), int(height * 0.08), "内容结构", font_size=22, color="0F172A", bold=True)
    _add_textbox(slide, int(width * 0.07), int(height * 0.15), int(width * 0.80), int(height * 0.07), title, font_size=11, color="64748B")

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(width * 0.07), int(height * 0.24), int(width * 0.86), int(height * 0.60))
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb("FFFFFF")
    card.line.color.rgb = _rgb("D7E3F4")

    outline_box = slide.shapes.add_textbox(int(width * 0.11), int(height * 0.30), int(width * 0.78), int(height * 0.48))
    frame = outline_box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, section in enumerate(sections[:MAX_OUTLINE_ITEMS], start=1):
        paragraph = frame.paragraphs[0] if index == 1 else frame.add_paragraph()
        paragraph.text = f"{index:02d}  {section.title}"
        paragraph.font.name = FONT_FAMILY
        paragraph.font.size = Pt(18)
        paragraph.font.bold = True
        paragraph.font.color.rgb = _rgb("1E293B")
        paragraph.space_after = Pt(10)


def _svg_dimensions(image_path: Path) -> tuple[int, int] | None:
    try:
        text = image_path.read_text(encoding="utf-8", errors="replace")
    except OSError:
        return None

    width_match = re.search(r'\bwidth=["\']([0-9.]+)(?:px|pt)?["\']', text)
    height_match = re.search(r'\bheight=["\']([0-9.]+)(?:px|pt)?["\']', text)
    if width_match and height_match:
        width = float(width_match.group(1))
        height = float(height_match.group(1))
    else:
        view_box = re.search(
            r'\bviewBox=["\'][-0-9.]+\s+[-0-9.]+\s+([0-9.]+)\s+([0-9.]+)["\']',
            text,
        )
        if not view_box:
            return None
        width = float(view_box.group(1))
        height = float(view_box.group(2))

    if width <= 0 or height <= 0:
        return None
    return int(round(width)), int(round(height))


def _image_dimensions(image_path: Path) -> tuple[int, int] | None:
    if image_path.suffix.lower() == ".svg":
        return _svg_dimensions(image_path)
    try:
        with Image.open(image_path) as image:
            return image.size
    except Exception:
        return None


def _fit_image_size(image_path: Path, max_width: int, max_height: int) -> tuple[int, int]:
    size = _image_dimensions(image_path)
    if size is None:
        return max_width, max_height
    width, height = size

    if width <= 0 or height <= 0:
        return max_width, max_height

    scale = min(max_width / width, max_height / height)
    scale = max(scale, 0.01)
    return max(int(width * scale), 1), max(int(height * scale), 1)


def _image_aspect_ratio(image_path: Path) -> float | None:
    size = _image_dimensions(image_path)
    if size is None:
        return None
    width, height = size
    if width <= 0 or height <= 0:
        return None
    return width / height


def _first_image_aspect_ratio(image_paths: list[Path] | None) -> float | None:
    for image_path in image_paths or []:
        if not image_path.exists():
            continue
        ratio = _image_aspect_ratio(image_path)
        if ratio is not None:
            return ratio
    return None


def _emu_to_pixels(length: int) -> int:
    return max(int(round(length / EMU_PER_PIXEL)), 1)


def _renderable_picture_path(image_path: Path, width: int, height: int) -> Path | None:
    if image_path.suffix.lower() != ".svg":
        return image_path
    if _CONVERT_SVG_TO_PNG_CACHED is None:
        return None

    width_px = _emu_to_pixels(width)
    height_px = _emu_to_pixels(height)
    cache_dir = image_path.parent / ".cache" / "simple_builder_svg_png"
    png_path = cache_dir / f"{image_path.stem}_{width_px}x{height_px}.png"
    ok = _CONVERT_SVG_TO_PNG_CACHED(image_path, png_path, width_px, height_px, cache_dir)
    if ok and png_path.exists():
        return png_path
    return None


def _add_picture_fit(slide, image_path: Path, left: int, top: int, width: int, height: int) -> bool:
    try:
        fitted_width, fitted_height = _fit_image_size(image_path, width, height)
        fitted_left = left + max((width - fitted_width) // 2, 0)
        fitted_top = top + max((height - fitted_height) // 2, 0)
        picture_path = _renderable_picture_path(image_path, fitted_width, fitted_height)
        if picture_path is None:
            return False
        slide.shapes.add_picture(str(picture_path), fitted_left, fitted_top, width=fitted_width, height=fitted_height)
        return True
    except Exception:
        return False


def _add_image_panel(slide, left: int, top: int, width: int, height: int, image_paths: list[Path]) -> int:
    valid_images = [image_path for image_path in image_paths[:2] if image_path.exists()]
    if not valid_images:
        return 0

    gap = int(height * 0.04)
    slot_count = len(valid_images)
    slot_height = height if slot_count == 1 else max((height - gap) // 2, 1)
    placed_count = 0

    for index, image_path in enumerate(valid_images):
        slot_top = top if slot_count == 1 else top + index * (slot_height + gap)
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, slot_top, width, slot_height)
        frame.fill.solid()
        frame.fill.fore_color.rgb = _rgb("F8FAFC")
        frame.line.color.rgb = _rgb("D7E3F4")

        inset = max(int(width * 0.04), int(height * 0.02))
        inner_left = left + inset
        inner_top = slot_top + inset
        inner_width = max(width - inset * 2, 1)
        inner_height = max(slot_height - inset * 2, 1)
        if _add_picture_fit(slide, image_path, inner_left, inner_top, inner_width, inner_height):
            placed_count += 1

    return placed_count


def _add_insight_layout(
    slide,
    panel_left: int,
    panel_top: int,
    panel_width: int,
    panel_height: int,
    bullets: list[str],
    lead: str,
    image_paths: list[Path] | None,
    profile: dict[str, float],
) -> int:
    gap = int(panel_width * profile["gap_ratio"])
    lead_height = int(panel_height * profile["lead_ratio"]) if lead else 0

    if lead_height:
        _add_lead_band(
            slide,
            panel_left + gap,
            panel_top + gap,
            panel_width - gap * 2,
            lead_height,
            lead,
            font_size=int(profile["lead_size"]),
        )

    content_top = panel_top + gap + lead_height + (gap if lead_height else 0)
    content_height = panel_height - (content_top - panel_top) - gap
    has_images = bool(image_paths)
    if has_images:
        content_left = panel_left + gap
        content_width = panel_width - gap * 2
        image_ratio = _first_image_aspect_ratio(image_paths)
        min_text_height = max(int(panel_height * 0.22), 90)
        min_text_width = max(int(panel_width * 0.28), 240)

        if image_ratio is not None and image_ratio >= 1.5:
            image_height = min(
                max(int(content_width / image_ratio), int(content_height * 0.30)),
                max(content_height - min_text_height - gap, 1),
            )
            if content_height - image_height - gap >= min_text_height:
                placed_count = _add_image_panel(
                    slide,
                    content_left,
                    content_top,
                    content_width,
                    image_height,
                    image_paths or [],
                )
                _add_bullet_box(
                    slide,
                    content_left,
                    content_top + image_height + gap,
                    content_width,
                    content_height - image_height - gap,
                    bullets,
                    font_size=int(profile["body_size"]),
                    color="1E293B",
                )
                return placed_count

        if image_ratio is None:
            image_width = int(content_width * 0.34)
        else:
            image_width = int(content_height * image_ratio)
        image_width = max(image_width, int(content_width * 0.22))
        image_width = min(image_width, int(content_width * 0.56))
        if content_width - image_width - gap < min_text_width:
            image_width = max(content_width - min_text_width - gap, int(content_width * 0.22))
        text_left = content_left
        text_width = max(content_width - image_width - gap, min_text_width)
        image_left = text_left + text_width + gap
        _add_bullet_box(
            slide,
            text_left,
            content_top,
            text_width,
            content_height,
            bullets,
            font_size=int(profile["body_size"]),
            color="1E293B",
        )
        return _add_image_panel(slide, image_left, content_top, image_width, content_height, image_paths or [])

    _add_bullet_box(
        slide,
        panel_left + gap,
        content_top,
        panel_width - gap * 2,
        content_height,
        bullets,
        font_size=int(profile["body_size"]),
        color="1E293B",
    )
    return 0


def _add_cards_layout(
    slide,
    panel_left: int,
    panel_top: int,
    panel_width: int,
    panel_height: int,
    bullets: list[str],
    lead: str,
    profile: dict[str, float],
) -> int:
    gap = int(panel_width * profile["gap_ratio"])
    lead_height = int(panel_height * profile["lead_ratio"]) if lead else 0
    if lead_height:
        _add_lead_band(
            slide,
            panel_left + gap,
            panel_top + gap,
            panel_width - gap * 2,
            lead_height,
            lead,
            font_size=int(profile["lead_size"]),
        )

    cards_top = panel_top + gap + lead_height + (gap if lead_height else 0)
    cards_height = panel_height - (cards_top - panel_top) - gap
    columns = 2 if len(bullets) > 1 else 1
    rows = max((len(bullets) + columns - 1) // columns, 1)
    card_width = max((panel_width - gap * (columns + 1)) // columns, 1)
    card_height = max((cards_height - gap * (rows + 1)) // rows, 1)

    for index, bullet in enumerate(bullets[:4]):
        row = index // columns
        column = index % columns
        card_left = panel_left + gap + column * (card_width + gap)
        card_top = cards_top + gap + row * (card_height + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb("F8FAFC")
        card.line.color.rgb = _rgb("BFDBFE")

        badge_size = min(int(card_width * 0.16), int(card_height * 0.28))
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, card_left + gap, card_top + gap, badge_size, badge_size)
        badge.fill.solid()
        badge.fill.fore_color.rgb = _rgb("2563EB")
        badge.line.fill.background()
        _add_textbox(
            slide,
            card_left + gap,
            card_top + gap + int(badge_size * 0.08),
            badge_size,
            badge_size,
            str(index + 1),
            font_size=11,
            color="FFFFFF",
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        _add_textbox(
            slide,
            card_left + gap,
            card_top + gap + badge_size + int(card_height * 0.06),
            card_width - gap * 2,
            card_height - badge_size - gap * 2,
            bullet,
            font_size=int(profile["body_size"]),
            color="1E293B",
        )
    return 0


def _add_comparison_layout(
    slide,
    panel_left: int,
    panel_top: int,
    panel_width: int,
    panel_height: int,
    bullets: list[str],
    lead: str,
    profile: dict[str, float],
) -> int:
    gap = int(panel_width * profile["gap_ratio"])
    lead_height = int(panel_height * profile["lead_ratio"]) if lead else 0
    if lead_height:
        _add_lead_band(
            slide,
            panel_left + gap,
            panel_top + gap,
            panel_width - gap * 2,
            lead_height,
            lead,
            font_size=int(profile["lead_size"]),
        )

    content_top = panel_top + gap + lead_height + (gap if lead_height else 0)
    content_height = panel_height - (content_top - panel_top) - gap
    column_width = max((panel_width - gap * 3) // 2, 1)
    left_items = bullets[: max((len(bullets) + 1) // 2, 1)]
    right_items = bullets[len(left_items):] or bullets[-1:]

    for offset, label, items in ((0, "视角一", left_items), (1, "视角二", right_items)):
        column_left = panel_left + gap + offset * (column_width + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, column_left, content_top, column_width, content_height)
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb("FFFFFF")
        card.line.color.rgb = _rgb("BFDBFE")
        _add_textbox(
            slide,
            column_left + gap,
            content_top + int(content_height * 0.06),
            column_width - gap * 2,
            int(content_height * 0.10),
            label,
            font_size=12,
            color="2563EB",
            bold=True,
        )
        _add_bullet_box(
            slide,
            column_left + gap,
            content_top + int(content_height * 0.18),
            column_width - gap * 2,
            int(content_height * 0.72),
            items,
            font_size=int(profile["body_size"]),
            color="1E293B",
        )
    return 0


def _add_timeline_layout(
    slide,
    panel_left: int,
    panel_top: int,
    panel_width: int,
    panel_height: int,
    bullets: list[str],
    lead: str,
    profile: dict[str, float],
) -> int:
    gap = int(panel_width * profile["gap_ratio"])
    lead_height = int(panel_height * profile["lead_ratio"]) if lead else 0
    if lead_height:
        _add_lead_band(
            slide,
            panel_left + gap,
            panel_top + gap,
            panel_width - gap * 2,
            lead_height,
            lead,
            font_size=int(profile["lead_size"]),
        )

    content_top = panel_top + gap + lead_height + (gap if lead_height else 0)
    content_height = panel_height - (content_top - panel_top) - gap
    line_left = panel_left + int(panel_width * 0.12)
    line_top = content_top + int(content_height * 0.08)
    line_height = int(content_height * 0.78)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_left, line_top, int(panel_width * 0.006), line_height)
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb("60A5FA")
    line.line.fill.background()

    step_gap = line_height // max(len(bullets), 1)
    for index, bullet in enumerate(bullets):
        node_top = line_top + index * step_gap
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, line_left - int(panel_width * 0.018), node_top, int(panel_width * 0.04), int(panel_width * 0.04))
        node.fill.solid()
        node.fill.fore_color.rgb = _rgb("2563EB")
        node.line.fill.background()
        _add_textbox(
            slide,
            line_left + int(panel_width * 0.05),
            node_top - int(content_height * 0.01),
            int(panel_width * 0.74),
            int(content_height * 0.14),
            bullet,
            font_size=int(profile["body_size"]),
            color="1E293B",
        )
    return 0


def _add_spotlight_layout(
    slide,
    panel_left: int,
    panel_top: int,
    panel_width: int,
    panel_height: int,
    bullets: list[str],
    lead: str,
    image_paths: list[Path] | None,
    profile: dict[str, float],
) -> int:
    gap = int(panel_width * profile["gap_ratio"])
    lead_text = lead or bullets[0]
    supporting_bullets = bullets[1:] if bullets and lead_text == bullets[0] and len(bullets) > 1 else bullets
    image_ratio = _first_image_aspect_ratio(image_paths)
    if image_paths and image_ratio is not None:
        available_image_height = panel_height - gap * 2
        image_width = int(available_image_height * image_ratio)
        image_width = max(image_width, int(panel_width * 0.24))
        image_width = min(image_width, int(panel_width * 0.52))
        lead_width = max(panel_width - image_width - gap * 3, int(panel_width * 0.34))
        image_width = panel_width - lead_width - gap * 3
    else:
        image_width = int(panel_width * 0.38)
        lead_width = int(panel_width * (0.42 if image_paths else 0.48))
    lead_height = int(panel_height * 0.54)
    _add_lead_band(
        slide,
        panel_left + gap,
        panel_top + gap,
        lead_width,
        lead_height,
        lead_text,
        font_size=int(profile["lead_size"] + 1),
        fill_color="DBEAFE",
        text_color="1D4ED8",
    )

    bullet_left = panel_left + gap
    bullet_top = panel_top + lead_height + gap * 2
    bullet_width = lead_width
    bullet_height = panel_height - lead_height - gap * 3
    _add_bullet_box(
        slide,
        bullet_left,
        bullet_top,
        bullet_width,
        bullet_height,
        supporting_bullets or bullets,
        font_size=int(profile["body_size"]),
        color="1E293B",
    )

    if image_paths:
        image_left = panel_left + lead_width + gap * 2
        return _add_image_panel(slide, image_left, panel_top + gap, image_width, panel_height - gap * 2, image_paths or [])
    return 0


def _add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    project_name: str,
    image_paths: list[Path] | None = None,
    *,
    layout_hint: str = "insight",
    page_rhythm: str = "anchor",
    lead: str = "",
) -> int:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    width = prs.slide_width
    height = prs.slide_height
    _add_background(slide, width, height, "F8FAFC")
    _add_slide_decor(slide, width, height, page_rhythm)

    top_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, int(height * 0.03))
    top_rule.fill.solid()
    top_rule.fill.fore_color.rgb = _rgb("2563EB")
    top_rule.line.fill.background()

    profile = _rhythm_profile(page_rhythm, height)
    resolved_layout = _resolve_layout_hint(layout_hint, image_paths, bullets)

    _add_textbox(
        slide,
        int(width * 0.07),
        int(height * 0.08),
        int(width * 0.74),
        int(height * 0.08),
        title,
        font_size=int(profile["title_size"]),
        color="0F172A",
        bold=True,
    )
    _add_textbox(slide, int(width * 0.07), int(height * 0.16), int(width * 0.50), int(height * 0.05), project_name, font_size=10, color="64748B")

    panel_left = int(width * 0.07)
    panel_top = int(height * 0.24)
    panel_width = int(width * 0.86)
    panel_height = int(height * 0.60)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, panel_left, panel_top, panel_width, panel_height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = _rgb("FFFFFF")
    panel.line.color.rgb = _rgb("D7E3F4")

    if resolved_layout == "cards":
        return _add_cards_layout(slide, panel_left, panel_top, panel_width, panel_height, bullets, lead, profile)
    if resolved_layout == "comparison":
        return _add_comparison_layout(slide, panel_left, panel_top, panel_width, panel_height, bullets, lead, profile)
    if resolved_layout == "timeline":
        return _add_timeline_layout(slide, panel_left, panel_top, panel_width, panel_height, bullets, lead, profile)
    if resolved_layout == "spotlight":
        return _add_spotlight_layout(slide, panel_left, panel_top, panel_width, panel_height, bullets, lead, image_paths, profile)
    if resolved_layout == "image_focus":
        return _add_spotlight_layout(slide, panel_left, panel_top, panel_width, panel_height, bullets, lead, image_paths, profile)
    return _add_insight_layout(slide, panel_left, panel_top, panel_width, panel_height, bullets, lead, image_paths, profile)


def _add_closing_slide(prs: Presentation, project_name: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    width = prs.slide_width
    height = prs.slide_height
    _add_background(slide, width, height, "0B1120")
    _add_textbox(slide, int(width * 0.10), int(height * 0.30), int(width * 0.70), int(height * 0.14), "已生成 PPT", font_size=30, color="F8FAFC", bold=True)
    _add_textbox(slide, int(width * 0.10), int(height * 0.48), int(width * 0.74), int(height * 0.10), "可直接下载并在 PowerPoint 中继续编辑。", font_size=16, color="CBD5E1")
    _add_textbox(slide, int(width * 0.10), int(height * 0.78), int(width * 0.74), int(height * 0.08), project_name, font_size=11, color="94A3B8")
